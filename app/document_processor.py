"""
Document Processor Module
Handles extraction of text from PDF, PPT/PPTX, and images (via Gemini OCR).
"""
from __future__ import annotations

import logging
import base64
from io import BytesIO
from threading import Lock
from typing import List, Dict, Any, Optional

import pypdf
from pptx import Presentation
from PIL import Image

from app.settings import AppSettings

logger = logging.getLogger("app.document_processor")

settings = AppSettings()

# Gemini client singleton
_CLIENT_LOCK = Lock()
_GEMINI_CLIENT = None


def _get_gemini_client():
    """Get or create Gemini client singleton."""
    global _GEMINI_CLIENT
    if _GEMINI_CLIENT is None:
        with _CLIENT_LOCK:
            if _GEMINI_CLIENT is None:
                if not settings.google_api_key:
                    raise RuntimeError("GOOGLE_API_KEY is not configured.")
                import google.genai as genai
                _GEMINI_CLIENT = genai.Client(api_key=settings.google_api_key)
    return _GEMINI_CLIENT


# =============================================================================
# Document Type Detection
# =============================================================================

SUPPORTED_DOC_EXTENSIONS = {
    ".pdf": "pdf",
    ".ppt": "ppt",
    ".pptx": "pptx",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
}

SKIP_EXTENSIONS = {".xlsx", ".xls", ".csv"}


def get_document_type(filename: str) -> str:
    """
    Determine document type from filename extension.
    
    Returns:
        'pdf', 'ppt', 'pptx', 'image', or 'unsupported'
    """
    ext = "." + filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    
    if ext in SKIP_EXTENSIONS:
        return "unsupported"
    
    return SUPPORTED_DOC_EXTENSIONS.get(ext, "unsupported")


def is_supported_document(filename: str) -> bool:
    """Check if file is a supported document type (non-Excel)."""
    return get_document_type(filename) != "unsupported"


# =============================================================================
# PDF Text Extraction
# =============================================================================

def extract_text_from_pdf(pdf_bytes: bytes) -> str:
    """
    Extract text from PDF bytes.
    
    Args:
        pdf_bytes: Raw PDF file content
        
    Returns:
        Extracted text as string
    """
    try:
        reader = pypdf.PdfReader(BytesIO(pdf_bytes))
        text_parts = []
        
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text.strip())
        
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.warning(f"Failed to extract text from PDF: {e}")
        return ""


# =============================================================================
# PPT/PPTX Text Extraction
# =============================================================================

def extract_text_from_pptx(pptx_bytes: bytes) -> str:
    """
    Extract text from PPT/PPTX bytes.
    
    Args:
        pptx_bytes: Raw PowerPoint file content
        
    Returns:
        Extracted text as string
    """
    try:
        prs = Presentation(BytesIO(pptx_bytes))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_texts.append(text)
            
            if slide_texts:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.warning(f"Failed to extract text from PPTX: {e}")
        return ""


# =============================================================================
# Image OCR (using Gemini LLM)
# =============================================================================

def extract_text_from_image(image_bytes: bytes, filename: str) -> str:
    """
    Extract text from image using Gemini LLM for OCR.
    
    Args:
        image_bytes: Raw image file content
        filename: Original filename (for format detection)
        
    Returns:
        Extracted/OCR'd text as string
    """
    try:
        client = _get_gemini_client()
        
        # Determine MIME type
        ext = filename.lower().rsplit(".", 1)[-1]
        mime_types = {
            "png": "image/png",
            "jpg": "image/jpeg",
            "jpeg": "image/jpeg",
        }
        mime_type = mime_types.get(ext, "image/png")
        
        # Encode image as base64
        image_b64 = base64.b64encode(image_bytes).decode("utf-8")
        
        # Create prompt for OCR
        prompt = """Ekstrak semua teks yang terlihat dalam gambar ini. 
Berikan hanya teks yang ditemukan, tanpa penjelasan tambahan.
Jika tidak ada teks yang terlihat, tulis "Tidak ada teks yang terdeteksi."
Jika gambar berisi tabel atau data terstruktur, format dalam bentuk yang rapi."""

        from google.genai import types as genai_types
        
        # Create content with image
        response = client.models.generate_content(
            model=settings.gemini_llm_model,
            contents=[
                genai_types.Content(
                    role="user",
                    parts=[
                        genai_types.Part(text=prompt),
                        genai_types.Part(
                            inline_data=genai_types.Blob(
                                mime_type=mime_type,
                                data=image_bytes
                            )
                        )
                    ]
                )
            ]
        )
        
        if response.text:
            return response.text.strip()
        return ""
        
    except Exception as e:
        logger.warning(f"Failed to OCR image {filename}: {e}")
        return ""


# =============================================================================
# Smart Semantic Text Chunking
# =============================================================================

def _split_into_semantic_blocks(text: str) -> List[str]:
    """
    Split text into semantic blocks (paragraphs, slide sections, etc).
    Never truncates - preserves complete blocks.
    """
    blocks = []
    
    # First, split by slide markers (for PPTX)
    if "--- Slide" in text:
        slide_parts = text.split("--- Slide")
        for i, part in enumerate(slide_parts):
            if i == 0 and not part.strip():
                continue
            if i > 0:
                part = "--- Slide" + part
            blocks.append(part.strip())
        return blocks
    
    # For other documents, split by double newlines (paragraphs)
    paragraphs = text.split("\n\n")
    for para in paragraphs:
        para = para.strip()
        if para:
            blocks.append(para)
    
    # If no paragraphs found, split by single newlines
    if not blocks:
        lines = text.split("\n")
        for line in lines:
            line = line.strip()
            if line:
                blocks.append(line)
    
    # If still nothing, return the whole text
    if not blocks and text.strip():
        blocks = [text.strip()]
    
    return blocks


def chunk_text(
    text: str,
    chunk_size: int = 1200,
    overlap: int = 0
) -> List[str]:
    """
    Smart semantic chunking that groups related content together.
    
    Strategy:
    1. Split text into semantic blocks (paragraphs, slides)
    2. Group blocks together up to chunk_size
    3. Never truncate a block - if a single block exceeds chunk_size, 
       keep it whole rather than cutting mid-sentence
    
    Args:
        text: Text to split
        chunk_size: Target chunk size in characters (soft limit)
        overlap: Ignored - kept for API compatibility
        
    Returns:
        List of text chunks with complete semantic units
    """
    if not text or not text.strip():
        return []
    
    text = text.strip()
    
    # If text is short enough, return as single chunk
    if len(text) <= chunk_size:
        return [text]
    
    # Split into semantic blocks
    blocks = _split_into_semantic_blocks(text)
    
    if not blocks:
        return [text]
    
    # Group blocks into chunks
    chunks = []
    current_chunk = []
    current_size = 0
    
    for block in blocks:
        block_size = len(block)
        
        # If single block exceeds chunk_size, add it as its own chunk
        # (never truncate)
        if block_size > chunk_size:
            # First, save current chunk if any
            if current_chunk:
                chunks.append("\n\n".join(current_chunk))
                current_chunk = []
                current_size = 0
            # Add large block as its own chunk
            chunks.append(block)
            continue
        
        # If adding this block would exceed size, start new chunk
        if current_size + block_size + 2 > chunk_size and current_chunk:
            chunks.append("\n\n".join(current_chunk))
            current_chunk = []
            current_size = 0
        
        current_chunk.append(block)
        current_size += block_size + 2  # +2 for "\n\n" separator
    
    # Don't forget the last chunk
    if current_chunk:
        chunks.append("\n\n".join(current_chunk))
    
    return chunks


def chunk_text_by_slides(text: str) -> List[str]:
    """
    Chunk PPTX text by individual slides.
    Each slide becomes its own chunk.
    """
    if "--- Slide" not in text:
        return chunk_text(text)
    
    chunks = []
    slide_parts = text.split("--- Slide")
    
    for i, part in enumerate(slide_parts):
        if i == 0 and not part.strip():
            continue
        slide_text = ("--- Slide" + part).strip() if i > 0 else part.strip()
        if slide_text:
            chunks.append(slide_text)
    
    return chunks if chunks else [text]


# =============================================================================
# Main Process Document Entry Point
# =============================================================================

def process_document(
    file_bytes: bytes,
    filename: str,
    chunk_size: int = 800,
    chunk_overlap: int = 100
) -> Optional[Dict[str, Any]]:
    """
    Process a document: extract text and create chunks.
    
    Args:
        file_bytes: Raw file content
        filename: Original filename
        chunk_size: Target chunk size
        chunk_overlap: Overlap between chunks
        
    Returns:
        Dict with filename, doc_type, full_text, and chunks
        None if unsupported document type
    """
    doc_type = get_document_type(filename)
    
    if doc_type == "unsupported":
        logger.info(f"Skipping unsupported document type: {filename}")
        return {"error": f"Unsupported document type: {filename}"}
    
    # Extract text based on type
    if doc_type == "pdf":
        full_text = extract_text_from_pdf(file_bytes)
    elif doc_type in ("ppt", "pptx"):
        full_text = extract_text_from_pptx(file_bytes)
    elif doc_type == "image":
        full_text = extract_text_from_image(file_bytes, filename)
    else:
        return {"error": f"Unknown document type: {doc_type}"}
    
    if not full_text:
        logger.warning(f"No text extracted from {filename}")
        return {
            "filename": filename,
            "doc_type": doc_type,
            "full_text": "",
            "chunks": [],
            "warning": "No text could be extracted"
        }
    
    # Chunk the text
    chunks = chunk_text(full_text, chunk_size, chunk_overlap)
    
    return {
        "filename": filename,
        "doc_type": doc_type,
        "full_text": full_text,
        "chunks": chunks,
        "chunk_count": len(chunks),
        "char_count": len(full_text)
    }
